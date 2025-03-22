import os
import uuid
from pathlib import Path
from pdf2docx import Converter as DocxConverter
from pdf2image import convert_from_path
import tabula
import pandas as pd
from PyPDF2 import PdfReader
from pptx import Presentation
from pptx.util import Inches
from PIL import Image


class PDFConverter:
    """
    Utility class for converting PDF files to various formats
    """
    
    @staticmethod
    def pdf_to_word(pdf_path, output_dir=None):
        """
        Convert PDF to Word document
        
        Args:
            pdf_path (str): Path to the PDF file
            output_dir (str, optional): Directory to save the converted file
            
        Returns:
            str: Path to the converted Word document
        """
        if output_dir is None:
            output_dir = os.path.join(os.getcwd(), 'converted', 'word')
            
        os.makedirs(output_dir, exist_ok=True)
        
        # Generate a unique filename
        output_filename = f"{uuid.uuid4()}.docx"
        output_path = os.path.join(output_dir, output_filename)
        
        # Convert PDF to Word
        try:
            cv = DocxConverter(pdf_path)
            cv.convert(output_path)
            cv.close()
            return output_path
        except Exception as e:
            raise Exception(f"Error converting PDF to Word: {str(e)}")
    
    @staticmethod
    def pdf_to_jpg(pdf_path, output_dir=None, dpi=300):
        """
        Convert PDF to JPG images
        
        Args:
            pdf_path (str): Path to the PDF file
            output_dir (str, optional): Directory to save the converted files
            dpi (int, optional): DPI for the output images
            
        Returns:
            list: List of paths to the converted JPG images
        """
        if output_dir is None:
            output_dir = os.path.join(os.getcwd(), 'converted', 'jpg')
            
        os.makedirs(output_dir, exist_ok=True)
        
        # Generate a unique folder for this conversion
        folder_name = str(uuid.uuid4())
        output_folder = os.path.join(output_dir, folder_name)
        os.makedirs(output_folder, exist_ok=True)
        
        # Convert PDF to images
        try:
            images = convert_from_path(pdf_path, dpi=dpi)
            image_paths = []
            
            for i, image in enumerate(images):
                image_path = os.path.join(output_folder, f"page_{i+1}.jpg")
                image.save(image_path, "JPEG")
                image_paths.append(image_path)
                
            return image_paths
        except Exception as e:
            raise Exception(f"Error converting PDF to JPG: {str(e)}")
    
    @staticmethod
    def pdf_to_excel(pdf_path, output_dir=None):
        """
        Convert PDF to Excel spreadsheet
        
        Args:
            pdf_path (str): Path to the PDF file
            output_dir (str, optional): Directory to save the converted file
            
        Returns:
            str: Path to the converted Excel file
        """
        if output_dir is None:
            output_dir = os.path.join(os.getcwd(), 'converted', 'excel')
            
        os.makedirs(output_dir, exist_ok=True)
        
        # Generate a unique filename
        output_filename = f"{uuid.uuid4()}.xlsx"
        output_path = os.path.join(output_dir, output_filename)
        
        # Convert PDF to Excel
        try:
            # Read tables from PDF
            tables = tabula.read_pdf(pdf_path, pages='all', multiple_tables=True)
            
            # Create Excel writer
            with pd.ExcelWriter(output_path) as writer:
                for i, table in enumerate(tables):
                    table.to_excel(writer, sheet_name=f'Sheet{i+1}', index=False)
                    
            return output_path
        except Exception as e:
            raise Exception(f"Error converting PDF to Excel: {str(e)}")
    
    @staticmethod
    def pdf_to_ppt(pdf_path, output_dir=None):
        """
        Convert PDF to PowerPoint presentation
        
        Args:
            pdf_path (str): Path to the PDF file
            output_dir (str, optional): Directory to save the converted file
            
        Returns:
            str: Path to the converted PowerPoint file
        """
        if output_dir is None:
            output_dir = os.path.join(os.getcwd(), 'converted', 'ppt')
            
        os.makedirs(output_dir, exist_ok=True)
        
        # Generate a unique filename
        output_filename = f"{uuid.uuid4()}.pptx"
        output_path = os.path.join(output_dir, output_filename)
        
        # Convert PDF to PowerPoint
        try:
            # First convert PDF pages to images
            image_paths = PDFConverter.pdf_to_jpg(pdf_path, dpi=200)
            
            # Create a PowerPoint presentation
            prs = Presentation()
            
            # Add each image as a slide
            for image_path in image_paths:
                slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
                
                # Add the image to the slide
                img = Image.open(image_path)
                width, height = img.size
                
                # Calculate aspect ratio
                aspect_ratio = width / height
                
                # Set slide dimensions (adjust as needed)
                if aspect_ratio > 1:  # Landscape
                    slide_width = Inches(10)
                    slide_height = slide_width / aspect_ratio
                else:  # Portrait
                    slide_height = Inches(7.5)
                    slide_width = slide_height * aspect_ratio
                
                # Add picture to slide
                slide.shapes.add_picture(image_path, Inches(0), Inches(0), width=slide_width, height=slide_height)
            
            # Save the presentation
            prs.save(output_path)
            
            return output_path
        except Exception as e:
            raise Exception(f"Error converting PDF to PowerPoint: {str(e)}")

    @staticmethod
    def word_to_pdf(word_path, output_dir=None):
        """
        Convert Word document to PDF
        
        Args:
            word_path (str): Path to the Word document
            output_dir (str, optional): Directory to save the converted file
            
        Returns:
            str: Path to the converted PDF file
        """
        if output_dir is None:
            output_dir = os.path.join(os.getcwd(), 'converted', 'pdf')
            
        os.makedirs(output_dir, exist_ok=True)
        
        # Generate a unique filename
        output_filename = f"{uuid.uuid4()}.pdf"
        output_path = os.path.join(output_dir, output_filename)
        
        # Convert Word to PDF
        try:
            from docx2pdf import convert
            convert(word_path, output_path)
            return output_path
        except Exception as e:
            raise Exception(f"Error converting Word to PDF: {str(e)}")
    
    @staticmethod
    def image_to_pdf(image_path, output_dir=None):
        """
        Convert image (JPG/PNG) to PDF
        
        Args:
            image_path (str): Path to the image file
            output_dir (str, optional): Directory to save the converted file
            
        Returns:
            str: Path to the converted PDF file
        """
        if output_dir is None:
            output_dir = os.path.join(os.getcwd(), 'converted', 'pdf')
            
        os.makedirs(output_dir, exist_ok=True)
        
        # Generate a unique filename
        output_filename = f"{uuid.uuid4()}.pdf"
        output_path = os.path.join(output_dir, output_filename)
        
        # Convert image to PDF
        try:
            from PIL import Image
            image = Image.open(image_path)
            image.save(output_path, "PDF", resolution=100.0)
            return output_path
        except Exception as e:
            raise Exception(f"Error converting image to PDF: {str(e)}")
    
    @staticmethod
    def excel_to_pdf(excel_path, output_dir=None):
        """
        Convert Excel spreadsheet to PDF
        
        Args:
            excel_path (str): Path to the Excel file
            output_dir (str, optional): Directory to save the converted file
            
        Returns:
            str: Path to the converted PDF file
        """
        if output_dir is None:
            output_dir = os.path.join(os.getcwd(), 'converted', 'pdf')
            
        os.makedirs(output_dir, exist_ok=True)
        
        # Generate a unique filename
        output_filename = f"{uuid.uuid4()}.pdf"
        output_path = os.path.join(output_dir, output_filename)
        
        # Convert Excel to PDF
        try:
            import win32com.client
            excel = win32com.client.Dispatch("Excel.Application")
            excel.Visible = False
            workbook = excel.Workbooks.Open(excel_path)
            workbook.ExportAsFixedFormat(0, output_path)
            workbook.Close()
            excel.Quit()
            return output_path
        except Exception as e:
            raise Exception(f"Error converting Excel to PDF: {str(e)}")
    
    @staticmethod
    def ppt_to_pdf(ppt_path, output_dir=None):
        """
        Convert PowerPoint presentation to PDF
        
        Args:
            ppt_path (str): Path to the PowerPoint file
            output_dir (str, optional): Directory to save the converted file
            
        Returns:
            str: Path to the converted PDF file
        """
        if output_dir is None:
            output_dir = os.path.join(os.getcwd(), 'converted', 'pdf')
            
        os.makedirs(output_dir, exist_ok=True)
        
        # Generate a unique filename
        output_filename = f"{uuid.uuid4()}.pdf"
        output_path = os.path.join(output_dir, output_filename)
        
        # Convert PowerPoint to PDF
        try:
            import win32com.client
            powerpoint = win32com.client.Dispatch("PowerPoint.Application")
            powerpoint.Visible = 1
            presentation = powerpoint.Presentations.Open(ppt_path)
            presentation.SaveAs(output_path, 32)
            presentation.Close()
            powerpoint.Quit()
            return output_path
        except Exception as e:
            raise Exception(f"Error converting PowerPoint to PDF: {str(e)}")
